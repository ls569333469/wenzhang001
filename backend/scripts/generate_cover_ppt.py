"""
P31: 投研报告 PPT 封面自动生成
用 python-pptx 生成精美中文封面，用于推特配图
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path
from datetime import datetime


def hex_to_rgb(hex_color: str) -> RGBColor:
    """#RRGGBB → RGBColor"""
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# === 配色方案 ===
COLORS = {
    "bg_dark": "#0B1120",
    "bg_card": "#111827",
    "text_white": "#F1F5F9",
    "text_gray": "#94A3B8",
    "text_dim": "#64748B",
    "accent_purple": "#A78BFA",
    "accent_cyan": "#22D3EE",
    "accent_green": "#34D399",
    "accent_amber": "#FBBF24",
    "accent_red": "#F87171",
    "accent_blue": "#6366F1",
    "tag_purple_bg": "#2E1065",
    "tag_cyan_bg": "#083344",
    "tag_amber_bg": "#422006",
    "tag_green_bg": "#052E16",
    "tag_red_bg": "#450A0A",
    "tag_orange_bg": "#431407",
}


def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(1)):
    """添加圆角矩形"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(fill_color)

    if border_color:
        shape.line.color.rgb = hex_to_rgb(border_color)
        shape.line.width = border_width
    else:
        shape.line.fill.background()

    # 圆角半径
    shape.adjustments[0] = 0.08
    return shape


def add_text_box(slide, left, top, width, height, text, font_name="微软雅黑",
                 font_size=12, font_color="#F1F5F9", bold=False, alignment=PP_ALIGN.LEFT):
    """添加文本框"""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.color.rgb = hex_to_rgb(font_color)
    p.font.bold = bold
    p.alignment = alignment
    return txbox


def add_tag(slide, left, top, text, bg_color, text_color, width=None):
    """添加标签"""
    tag_w = width or Inches(1.3)
    tag_h = Inches(0.3)
    shape = add_rounded_rect(slide, left, top, tag_w, tag_h, bg_color, text_color)
    shape.adjustments[0] = 0.5  # 更大的圆角

    tf = shape.text_frame
    tf.word_wrap = False
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "微软雅黑"
    p.font.size = Pt(9)
    p.font.color.rgb = hex_to_rgb(text_color)
    p.font.bold = True
    # 垂直居中
    tf.paragraphs[0].space_before = Pt(1)

    return shape


def add_rank_circle(slide, left, top, number, color):
    """排名圆圈"""
    size = Inches(0.5)
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(COLORS["bg_dark"])
    shape.line.color.rgb = hex_to_rgb(color)
    shape.line.width = Pt(2.5)

    tf = shape.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = str(number)
    run.font.name = "Consolas"
    run.font.size = Pt(18)
    run.font.color.rgb = hex_to_rgb(color)
    run.font.bold = True

    return shape


def generate_cover(projects: list, date_str: str = None, output_path: str = None):
    """
    生成投研报告 PPT 封面

    projects: [
        {
            "name": "Reveel",
            "twitter": "@r3vl_xyz",
            "stage": "Pre-TGE",
            "category": "AI 支付 Infra",
            "event": "Binance Booster",
            "kol_24h": 20,
            "verdict": "观望",       # 观望/买入/回避
        },
        ...
    ]
    """
    if not date_str:
        date_str = datetime.now().strftime("%Y.%m.%d")

    prs = Presentation()

    # 16:9 宽屏
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局

    # === 1. 背景色 ===
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(COLORS["bg_dark"])

    # === 2. 装饰元素 ===
    # 顶部渐变装饰条
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(0.05)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = hex_to_rgb(COLORS["accent_blue"])
    top_bar.line.fill.background()

    # 左侧装饰竖线
    left_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.8), Inches(0.6),
        Inches(0.04), Inches(1.0)
    )
    left_line.fill.solid()
    left_line.fill.fore_color.rgb = hex_to_rgb(COLORS["accent_purple"])
    left_line.line.fill.background()

    # === 3. HEADER ===
    # 标题：Daily Alpha
    add_text_box(slide,
        Inches(1.1), Inches(0.55), Inches(5), Inches(0.8),
        "🌊 Daily Alpha", "微软雅黑", 36, COLORS["text_white"], True
    )

    # 副标题
    add_text_box(slide,
        Inches(1.1), Inches(1.2), Inches(6), Inches(0.4),
        "leak.me × Surf AI · KOL 关注追踪 + 深度投研",
        "微软雅黑", 13, COLORS["text_dim"], False
    )

    # 日期（右上角）
    add_text_box(slide,
        Inches(10.0), Inches(0.55), Inches(2.5), Inches(0.6),
        date_str, "Consolas", 28, COLORS["text_gray"], True, PP_ALIGN.RIGHT
    )

    # 数据来源
    add_text_box(slide,
        Inches(10.0), Inches(1.1), Inches(2.5), Inches(0.35),
        f"24h Trending · Top {len(projects)}",
        "微软雅黑", 12, COLORS["text_dim"], False, PP_ALIGN.RIGHT
    )

    # === 分隔线 ===
    divider = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1.1), Inches(1.75), Inches(11.4), Inches(0.01)
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = hex_to_rgb("#1E293B")
    divider.line.fill.background()

    # === 4. 项目卡片 ===
    card_top = Inches(2.1)
    card_height = Inches(1.3)
    card_gap = Inches(0.2)

    # 颜色映射
    rank_colors = [COLORS["accent_amber"], COLORS["accent_blue"], COLORS["accent_cyan"],
                   COLORS["accent_green"], COLORS["accent_purple"]]

    stage_styles = {
        "Pre-TGE":  (COLORS["tag_purple_bg"], COLORS["accent_purple"]),
        "概念期":    (COLORS["tag_orange_bg"], COLORS["accent_amber"]),
        "已上线":    (COLORS["tag_green_bg"], COLORS["accent_green"]),
        "已融资":    (COLORS["tag_cyan_bg"], COLORS["accent_cyan"]),
        "成熟协议":  (COLORS["tag_cyan_bg"], COLORS["accent_cyan"]),
    }

    verdict_styles = {
        "观望": (COLORS["tag_green_bg"], COLORS["accent_green"], "👀 观望"),
        "买入": (COLORS["tag_cyan_bg"], COLORS["accent_cyan"], "🚀 买入"),
        "回避": (COLORS["tag_red_bg"], COLORS["accent_red"], "⚠️ 回避"),
    }

    for i, proj in enumerate(projects):
        y = card_top + (card_height + card_gap) * i
        color_idx = i % len(rank_colors)
        rank_color = rank_colors[color_idx]

        # 卡片背景
        add_rounded_rect(slide,
            Inches(1.1), y, Inches(11.4), card_height,
            COLORS["bg_card"], "#1E293B"
        )

        # 排名圆圈
        add_rank_circle(slide, Inches(1.5), y + Inches(0.38), i + 1, rank_color)

        # 项目名称
        add_text_box(slide,
            Inches(2.2), y + Inches(0.18), Inches(3), Inches(0.45),
            proj["name"], "微软雅黑", 22, COLORS["text_white"], True
        )

        # Twitter handle
        add_text_box(slide,
            Inches(2.2), y + Inches(0.58), Inches(2), Inches(0.3),
            proj.get("twitter", ""), "Consolas", 11, COLORS["text_dim"]
        )

        # 标签：阶段
        stage = proj.get("stage", "Pre-TGE")
        stage_bg, stage_color = stage_styles.get(stage, (COLORS["tag_purple_bg"], COLORS["accent_purple"]))
        tag_left = Inches(4.5)
        add_tag(slide, tag_left, y + Inches(0.45), stage, stage_bg, stage_color, Inches(1.1))

        # 标签：类别
        add_tag(slide, tag_left + Inches(1.2), y + Inches(0.45),
                proj.get("category", ""), COLORS["tag_cyan_bg"], COLORS["accent_cyan"], Inches(1.5))

        # 标签：关键事件
        event = proj.get("event", "")
        if event:
            add_tag(slide, tag_left + Inches(2.8), y + Inches(0.45),
                    event, COLORS["tag_amber_bg"], COLORS["accent_amber"], Inches(1.6))

        # KOL 数字（右侧）
        kol_num = proj.get("kol_24h", 0)
        kol_color = COLORS["accent_green"] if kol_num >= 15 else COLORS["accent_amber"]
        add_text_box(slide,
            Inches(9.2), y + Inches(0.15), Inches(1.5), Inches(0.55),
            f"+{kol_num}", "Consolas", 30, kol_color, True, PP_ALIGN.CENTER
        )
        add_text_box(slide,
            Inches(9.2), y + Inches(0.7), Inches(1.5), Inches(0.3),
            "24H KOL", "微软雅黑", 10, COLORS["text_dim"], True, PP_ALIGN.CENTER
        )

        # 投资判定（最右侧）
        verdict = proj.get("verdict", "观望")
        v_bg, v_color, v_text = verdict_styles.get(verdict, verdict_styles["观望"])
        add_tag(slide, Inches(10.8), y + Inches(0.45), v_text, v_bg, v_color, Inches(1.3))

    # === 5. FOOTER ===
    footer_y = Inches(6.8)

    # 免责声明
    add_text_box(slide,
        Inches(1.1), footer_y, Inches(8), Inches(0.3),
        "⚠️ 仅供研究参考，不构成投资建议 · Data: Surf AI + leak.me KOL Tracker",
        "微软雅黑", 10, COLORS["text_dim"]
    )

    # 品牌
    add_text_box(slide,
        Inches(10.5), footer_y, Inches(2), Inches(0.3),
        "Quantum Studio", "微软雅黑", 12, COLORS["text_gray"], True, PP_ALIGN.RIGHT
    )

    # === 保存 ===
    if not output_path:
        output_dir = Path(__file__).parent.parent.parent / "reports" / "research"
        output_dir.mkdir(parents=True, exist_ok=True)
        output_path = str(output_dir / f"daily_alpha_{date_str.replace('.', '')}.pptx")

    prs.save(output_path)
    return output_path


# === 主程序 ===
if __name__ == "__main__":
    # 使用今天的投研数据
    projects = [
        {
            "name": "Reveel",
            "twitter": "@r3vl_xyz",
            "stage": "Pre-TGE",
            "category": "AI 支付 Infra",
            "event": "Binance Booster",
            "kol_24h": 20,
            "verdict": "观望",
        },
        {
            "name": "Saturn Credit",
            "twitter": "@saturn_credit",
            "stage": "Pre-TGE",
            "category": "BTC 收益 DeFi",
            "event": "审计完成",
            "kol_24h": 20,
            "verdict": "观望",
        },
        {
            "name": "TechDollar",
            "twitter": "@techdollarhq",
            "stage": "概念期",
            "category": "DeFi 私人信贷",
            "event": "Waitlist 开启",
            "kol_24h": 20,
            "verdict": "回避",
        },
    ]

    output = generate_cover(projects, "2026.02.25")
    print(f"✅ PPT 封面已生成: {output}")
    print(f"   请用 PowerPoint 打开查看效果，截图即可用于推特配图")
